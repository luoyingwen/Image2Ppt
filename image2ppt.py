import os
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from displaypicture import ImageViewer
import logging

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# 加载环境变量
load_dotenv()

def crop_image(original_image_path, cropped_image_path, top_left, bottom_right):
    """
    裁剪图片并保存。
    
    Args:
        original_image_path (str): 原始图片路径
        cropped_image_path (str): 裁剪后图片保存路径  
        top_left (tuple): 左上角坐标 (x, y)
        bottom_right (tuple): 右下角坐标 (x, y)
    """
    try:
        image = Image.open(original_image_path)
        crop_area = (top_left[0], top_left[1], bottom_right[0], bottom_right[1])
        cropped_image = image.crop(crop_area)
        
        # 在保存之前删除已存在的文件
        if os.path.exists(cropped_image_path):
            os.remove(cropped_image_path)
        
        cropped_image.save(cropped_image_path)
        logger.info(f"成功裁剪并保存图片: {cropped_image_path}")
    except Exception as e:
        logger.error(f"裁剪图片时出错: {str(e)}")
        raise

def get_sorted_png_files(folder):
    """
    获取指定文件夹中的PNG文件，并按最后修改时间排序。
    
    Args:
        folder (str): 文件夹路径
        
    Returns:
        list: 排序后的PNG文件列表
    """
    try:
        png_files = [f for f in os.listdir(folder) if f.lower().endswith('.png')]
        
        if not png_files:
            logger.warning(f"文件夹 {folder} 中未找到PNG文件")
            return []

        logger.info(f"找到 {len(png_files)} 个PNG文件")
        logger.debug("排序前的文件名：\n%s", "\n".join(png_files))

        # 获取图像最后修改时间并排序
        png_files.sort(key=lambda x: os.stat(os.path.join(folder, x)).st_mtime)
        
        logger.debug("排序后的文件名：\n%s", "\n".join(png_files))
        return png_files
    except Exception as e:
        logger.error(f"获取PNG文件列表时出错: {str(e)}")
        raise

def create_ppt_from_png_files(png_files, output_pptx_path):
    """
    将PNG文件列表生成PPTX文件。
    
    Args:
        png_files (list): PNG文件路径列表
        output_pptx_path (str): 输出PPTX文件路径
    """
    try:
        presentation = Presentation()
        
        # 设置幻灯片大小为16:9
        presentation.slide_width = Inches(16)
        presentation.slide_height = Inches(9)

        for png_file in png_files:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            
            # 添加图片并使其填满幻灯片
            slide.shapes.add_picture(png_file, 0, 0, 
                                   width=presentation.slide_width, 
                                   height=presentation.slide_height)
            logger.info(f"已添加图片到幻灯片: {png_file}")

        # 保存PPTX文件
        presentation.save(output_pptx_path)
        logger.info(f"成功保存PPTX文件: {output_pptx_path}")
    except Exception as e:
        logger.error(f"创建PPT文件时出错: {str(e)}")
        raise

def crop_images_in_folder(input_folder, output_folder, top_left, bottom_right):
    """
    遍历文件夹中的PNG图片，按最后修改时间裁剪并保存。
    
    Args:
        input_folder (str): 输入文件夹路径
        output_folder (str): 输出文件夹路径
        top_left (tuple): 左上角坐标
        bottom_right (tuple): 右下角坐标
        
    Returns:
        list: 裁剪后的图片文件列表
    """
    try:
        png_files = get_sorted_png_files(input_folder)
        if not png_files:
            raise ValueError(f"在 {input_folder} 中未找到PNG文件")
            
        os.makedirs(output_folder, exist_ok=True)

        for png_file in png_files:
            original_image_path = os.path.join(input_folder, png_file)
            cropped_image_name = f"{os.path.splitext(png_file)[0]}_cropped.png"
            cropped_image_path = os.path.join(output_folder, cropped_image_name)
            logger.info(f"正在裁剪: {original_image_path} -> {cropped_image_path}")
            crop_image(original_image_path, cropped_image_path, top_left, bottom_right)

        return get_sorted_png_files(output_folder)
    except Exception as e:
        logger.error(f"批量裁剪图片时出错: {str(e)}")
        raise

def main():
    try:
        # 从环境变量读取参数
        input_folder = os.getenv('INPUT_FOLDER')
        output_folder = os.getenv('OUTPUT_FOLDER')
        
        if not input_folder or not output_folder:
            raise ValueError("请在.env文件中设置INPUT_FOLDER和OUTPUT_FOLDER")

        # 从input_folder中选择一张图片用于预览
        preview_image = None
        for file in os.listdir(input_folder):
            if file.lower().endswith(('.png', '.jpg', '.jpeg')):
                preview_image = os.path.join(input_folder, file)
                logger.info(f"选择预览图片: {preview_image}")
                break

        if preview_image is None:
            raise FileNotFoundError("在输入文件夹中未找到PNG或JPG格式的图片")

        # 使用ImageViewer显示预览图片并获取坐标
        viewer = ImageViewer(preview_image)
        coordinates = viewer.show()
        top_left = coordinates[0]
        bottom_right = coordinates[1]

        logger.info(f"选择的裁剪区域: 左上角{top_left}, 右下角{bottom_right}")

        # 执行裁剪
        cropped_images = crop_images_in_folder(input_folder, output_folder, top_left, bottom_right)

        # 生成PPTX文件
        output_pptx_path = os.path.join(input_folder, 'output_presentation.pptx')
        create_ppt_from_png_files([os.path.join(output_folder, img) for img in cropped_images], 
                                output_pptx_path)

        logger.info(f"处理完成! PPTX文件已保存为: {output_pptx_path}")
        
    except Exception as e:
        logger.error(f"程序执行出错: {str(e)}")
        raise

if __name__ == "__main__":
    main()
